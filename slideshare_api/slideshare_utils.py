
import os
import requests
from bs4 import BeautifulSoup
from fpdf import FPDF
from pptx import Presentation
from PIL import Image
import zipfile
from docx.shared import Inches
from docx import Document
from io import BytesIO
import concurrent.futures


def fetch_image_urls(slideshare_url):
    """
    Fetches image URLs from a SlideShare URL.
    """
    try:
        response = requests.get(slideshare_url)
        response.raise_for_status()
        html = response.content
        soup = BeautifulSoup(html, "html.parser")
        
        # Try different selectors for slide images
        images = soup.find_all("img", {"data-testid": "vertical-slide-image"})
        if not images:
            # Fallback to other possible selectors
            images = soup.find_all("img", class_=lambda x: x and "slide" in x.lower()) or \
                     soup.find_all("img", {"data-test": "slide-image"}) or \
                     soup.find_all("img", src=lambda x: x and "slide" in x.lower())
        
        image_urls = []

        for image in images:
            # Try to get image URL from srcset first, then fallback to src
            srcset = image.get("srcset")
            if srcset:
                try:
                    # Parse srcset to get the highest resolution image
                    image_url = srcset.split("w, ")[-1].split(" ")[0]
                except (IndexError, AttributeError):
                    image_url = image.get("src")
            else:
                image_url = image.get("src")
            
            # Only add valid image URLs
            if image_url and image_url.endswith((".jpg", ".jpeg", ".png", ".webp")):
                # Ensure absolute URL
                if image_url.startswith("//"):
                    image_url = "https:" + image_url
                elif image_url.startswith("/"):
                    from urllib.parse import urlparse
                    parsed_url = urlparse(slideshare_url)
                    image_url = f"{parsed_url.scheme}://{parsed_url.netloc}{image_url}"
                image_urls.append(image_url)

        return image_urls

    except requests.RequestException as e:
        raise Exception(f"Failed to fetch SlideShare page: {e}")
    except Exception as e:
        raise Exception(f"Failed to fetch image URLs: {e}")


def download_images(image_urls):
    """
    Downloads images from a list of URLs and returns their paths.
    """
    def fetch_and_compress(args):
        index, image_url = args
        try:
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            # Use BytesIO for in-memory processing
            img_bytes = BytesIO(response.content)
            with Image.open(img_bytes) as img:
                # Resize if too large (e.g., max width 800px)
                max_width = 800
                if img.width > max_width:
                    ratio = max_width / img.width
                    new_size = (max_width, int(img.height * ratio))
                    img = img.resize(new_size, Image.LANCZOS)
                # Always save as JPEG, quality 55 (smaller size, still readable)
                out_path = f"temp_image_{index}.jpg"
                img.convert("RGB").save(out_path, "JPEG", quality=55, optimize=True, progressive=True)
            return out_path
        except Exception as e:
            raise Exception(f"Failed to download or process image {image_url}: {e}")

    try:
        image_paths = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=8) as executor:
            results = list(executor.map(fetch_and_compress, enumerate(image_urls)))
        image_paths.extend(results)
        return image_paths
    except Exception as e:
        raise Exception(f"Failed to download images: {e}")


def convert_to_jpeg(image_path):
    """
    Converts an image to JPEG format if it's not already in JPEG.
    Returns the path of the converted JPEG file.
    """
    try:
        with Image.open(image_path) as img:
            if img.format != "JPEG":
                jpeg_path = f"{os.path.splitext(image_path)[0]}.jpg"
                img.convert("RGB").save(jpeg_path, "JPEG")
                return jpeg_path
            return image_path
    except Exception as e:
        raise Exception(f"Failed to convert image to JPEG: {e}")


def convert_images_to_pdf(image_paths):
    """
    Converts a list of image paths to a PDF.
    """
    try:
        pdf = FPDF()
        for image_path in image_paths:
            # Convert to JPEG if necessary
            jpeg_path = convert_to_jpeg(image_path)

            # Add the image to the PDF
            pdf.add_page()
            pdf.image(jpeg_path, x=10, y=10, w=190)

            # Cleanup the temporary JPEG file (if it was created)
            if jpeg_path != image_path:
                os.remove(jpeg_path)

            # Cleanup the original image file
            os.remove(image_path)

        pdf_path = "output.pdf"
        pdf.output(pdf_path)
        return pdf_path

    except Exception as e:
        raise Exception(f"Failed to convert images to PDF: {e}")


def convert_images_to_ppt(image_paths):
    """
    Converts a list of image paths to a PPT.
    """
    try:
        prs = Presentation()
        for image_path in image_paths:
            # Convert to JPEG if necessary
            jpeg_path = convert_to_jpeg(image_path)

            # Add the image to the PPT
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
            slide.shapes.add_picture(jpeg_path, 0, 0, width=prs.slide_width)

            # Cleanup the temporary JPEG file (if it was created)
            if jpeg_path != image_path:
                os.remove(jpeg_path)

            # Cleanup the original image file
            os.remove(image_path)

        ppt_path = "output.pptx"
        prs.save(ppt_path)
        return ppt_path

    except Exception as e:
        raise Exception(f"Failed to convert images to PPT: {e}")


def convert_images_to_word(image_paths):
    """
    Converts a list of image paths to a Word document.
    """
    try:
        doc = Document()
        for image_path in image_paths:
            # Convert to JPEG if necessary
            jpeg_path = convert_to_jpeg(image_path)

            # Add the image to the Word document
            doc.add_picture(jpeg_path, width=Inches(6))

            # Cleanup the temporary JPEG file (if it was created)
            if jpeg_path != image_path:
                os.remove(jpeg_path)

            # Cleanup the original image file
            os.remove(image_path)

        doc_path = "output.docx"
        doc.save(doc_path)
        return doc_path

    except Exception as e:
        raise Exception(f"Failed to convert images to Word: {e}")


def compress_file(file_path):
    """
    Compresses a file (PDF, PPT, or Word) into a ZIP archive.
    """
    try:
        zip_path = f"{file_path}.zip"
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            zipf.write(file_path, os.path.basename(file_path))

        os.remove(file_path)  # Cleanup
        return zip_path

    except Exception as e:
        raise Exception(f"Failed to compress file: {e}")
    

def optimize_ppt_size(ppt_path):
    """Reduce PPT file size by compressing images"""
    try:
        prs = Presentation(ppt_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    # Compress images in the PPT
                    shape.image.blob = compress_image(shape.image.blob)
        optimized_path = ppt_path.replace(".pptx", "_optimized.pptx")
        prs.save(optimized_path)
        return optimized_path
    except Exception as e:
        print(f"PPT optimization failed: {e}")
        return ppt_path  # Fallback to original file

def optimize_word_size(word_path):
    """Reduce Word file size by compressing images"""
    try:
        doc = Document(word_path)
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run.element.xpath(".//pic:pic"):
                    # Compress images in the Word document
                    for img in run.element.xpath(".//pic:pic"):
                        img_blob = img.xpath(".//a:blob")[0].text
                        img.xpath(".//a:blob")[0].text = compress_image(img_blob)
        optimized_path = word_path.replace(".docx", "_optimized.docx")
        doc.save(optimized_path)
        return optimized_path
    except Exception as e:
        print(f"Word optimization failed: {e}")
        return word_path  # Fallback to original file

def compress_image(image_data):
    """Compress image data using Pillow"""
    from PIL import Image
    from io import BytesIO
    try:
        img = Image.open(BytesIO(image_data))
        output = BytesIO()
        img.save(output, format="JPEG", quality=85)  # Adjust quality as needed
        return output.getvalue()
    except Exception as e:
        print(f"Image compression failed: {e}")
        return image_data  # Fallback to original image
    
def compress_file(file_path):
    """Compress a file into a ZIP archive"""
    zip_path = f"{file_path}.zip"
    try:
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            zipf.write(file_path, arcname=os.path.basename(file_path))
        return zip_path
    except Exception as e:
        print(f"Compression failed: {e}")
        return None